#!/usr/bin/env python3
"""Convert RTR presentation HTML files to PPTX format for Canva import."""

import os, re
from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ──────────────────────────────────
DARK     = RGBColor(0x0c, 0x11, 0x15)
NAVY     = RGBColor(0x1D, 0x47, 0x71)
MED_BLUE = RGBColor(0x03, 0x65, 0xb2)
LT_BLUE  = RGBColor(0x6d, 0xa3, 0xd9)
CORAL    = RGBColor(0xff, 0x5e, 0x5b)
WHITE    = RGBColor(0xff, 0xff, 0xff)
OFF_WHT  = RGBColor(0xf0, 0xf5, 0xfb)
GRAY_TXT = RGBColor(0x33, 0x4e, 0x68)
GHOST    = RGBColor(0xcc, 0xdb, 0xea)

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.65)


# ── Helpers ────────────────────────────────────────

def clean(text):
    return re.sub(r'\s+', ' ', (text or '').strip())

def get_text(el):
    return clean(el.get_text()) if el else ''

def add_rect(slide, x, y, w, h, fill=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def add_txb(slide, x, y, w, h, bg=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    else:
        txb.fill.background()
    txb.line.fill.background()
    return tf

def para(tf, text, size=16, bold=False, italic=False,
         color=DARK, align=PP_ALIGN.LEFT, font='Montserrat',
         space_before=Pt(0), space_after=Pt(5)):
    if not text:
        return
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return p


# ── Content extraction ────────────────────────────

def extract_items(el):
    """Walk an element tree and yield (type, text) pairs."""
    if not el or not isinstance(el, Tag):
        return
    classes = el.get('class', [])

    # Named structural patterns — handle and stop recursing into children
    if 'blist' in classes or (el.name == 'ul' and 'blist' in classes):
        for li in el.find_all('li', recursive=False):
            t = clean(li.get_text())
            if t:
                yield ('bullet', t)
        return

    if el.name == 'li':
        t = clean(el.get_text())
        if t:
            yield ('bullet', t)
        return

    if any(c in classes for c in ('card', 'ctrl-card', 'shift-box')):
        h = el.find(['h3', 'h4'])
        if h:
            yield ('card-head', clean(h.get_text()))
        for p in el.find_all('p', recursive=False):
            t = clean(p.get_text())
            if t:
                yield ('card-body', t)
        return

    if 'callout' in classes or 'callout-blue' in classes:
        label = el.find(class_='clabel')
        if label:
            yield ('callout-label', clean(label.get_text()))
        for p in el.find_all('p', recursive=False):
            t = clean(p.get_text())
            if t:
                yield ('callout-text', t)
        return

    if 'fc' in classes:
        label = el.find(class_='fc-label')
        if label:
            yield ('feature-label', clean(label.get_text()))
        for p in el.find_all('p', recursive=False):
            t = clean(p.get_text())
            if t:
                yield ('feature-text', t)
        return

    if 'hlevel' in classes:
        num  = el.find(class_='hl-num')
        name = el.find(class_='hl-name')
        tags = el.find_all(class_='hl-tag')
        label = ''
        if num:  label += clean(num.get_text()) + '. '
        if name: label += clean(name.get_text())
        if label:
            yield ('section-head', label)
        for t in tags:
            yield ('bullet', clean(t.get_text()))
        return

    if el.name in ('h2', 'h3') and el.parent and 'bd' not in el.parent.get('class', []):
        t = clean(el.get_text())
        if t:
            yield ('section-head', t)
        return

    if el.name == 'p':
        t = clean(el.get_text())
        if t:
            yield ('para', t)
        return

    # Recurse into generic divs / containers
    for child in el.children:
        if isinstance(child, Tag):
            yield from extract_items(child)


# ── Slide builders ────────────────────────────────

def build_title(prs, el):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=DARK)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=MED_BLUE)

    eyebrow  = el.find(class_='eyebrow')
    h1       = el.find('h1')
    subtitle = el.find(class_='subtitle')

    y = Inches(1.6)
    if eyebrow:
        tf = add_txb(slide, M, y, W - 2*M, Inches(0.45))
        para(tf, clean(eyebrow.get_text()), size=12, bold=True, color=LT_BLUE)
        y += Inches(0.5)

    if h1:
        tf = add_txb(slide, M, y, W - 2*M, Inches(2.2))
        para(tf, clean(h1.get_text()), size=54, bold=True, color=WHITE, font='Georgia')
        y += Inches(2.2)

    add_rect(slide, M, y, Inches(0.85), Inches(0.055), fill=MED_BLUE)
    y += Inches(0.28)

    if subtitle:
        tf = add_txb(slide, M, y, W - 2*M, Inches(0.65))
        para(tf, clean(subtitle.get_text()), size=22, color=RGBColor(0xaa, 0xbb, 0xcc))


def build_divider(prs, el):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=OFF_WHT)
    add_rect(slide, 0, 0, Inches(0.14), H, fill=NAVY)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=MED_BLUE)

    ghost = el.find(class_='ghost-num')
    h2    = el.find('h2')
    p     = el.find('p')

    if ghost:
        tf = add_txb(slide, Inches(7), Inches(0.3), Inches(5.5), Inches(3))
        para(tf, clean(ghost.get_text()), size=110, bold=True, color=GHOST)

    y = Inches(2.0)
    x = Inches(0.75)
    w = W - x - M

    if h2:
        tf = add_txb(slide, x, y, w, Inches(1.5))
        para(tf, clean(h2.get_text()), size=42, bold=True, color=NAVY, font='Georgia')
        y += Inches(1.55)

    if p:
        tf = add_txb(slide, x, y, w, Inches(2.0))
        para(tf, clean(p.get_text()), size=20, color=GRAY_TXT)


def build_nutshell(prs, el):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    LEFT_W = Inches(4.8)
    add_rect(slide, 0, 0, LEFT_W, H, fill=NAVY)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=MED_BLUE)

    ns_left  = el.find(class_='ns-left')
    ns_right = el.find(class_='ns-right')

    if ns_left:
        y = Inches(1.8)
        badge = ns_left.find(class_='badge')
        h2    = ns_left.find('h2')
        if badge:
            tf = add_txb(slide, M, y, LEFT_W - 2*M, Inches(0.45))
            para(tf, clean(badge.get_text()), size=11, bold=True, color=LT_BLUE)
            y += Inches(0.5)
        if h2:
            tf = add_txb(slide, M, y, LEFT_W - 2*M, Inches(2.5))
            para(tf, clean(h2.get_text()), size=34, bold=True, color=WHITE, font='Georgia')

    if ns_right:
        RX = LEFT_W + M
        RW = W - RX - M
        y = Inches(1.0)
        for item in ns_right.find_all(class_='ns-item'):
            act_label = item.find(class_='act-label')
            p         = item.find('p')
            if act_label:
                tf = add_txb(slide, RX, y, RW, Inches(0.4))
                para(tf, clean(act_label.get_text()), size=12, bold=True, color=NAVY)
                y += Inches(0.4)
            if p:
                tf = add_txb(slide, RX, y, RW, Inches(1.5))
                para(tf, clean(p.get_text()), size=17, color=GRAY_TXT)
                y += Inches(1.6)


def build_std(prs, el):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    add_rect(slide, 0, 0, W, Inches(1.15), fill=NAVY)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=MED_BLUE)

    hd = el.find(class_='hd')
    bd = el.find(class_='bd')

    if hd:
        tag = hd.find(class_='tag')
        h2  = hd.find('h2')
        if tag:
            tf = add_txb(slide, M, Inches(0.14), W - 2*M, Inches(0.38))
            para(tf, clean(tag.get_text()), size=11, bold=True, color=LT_BLUE)
        if h2:
            tf = add_txb(slide, M, Inches(0.46), W - 2*M, Inches(0.68))
            para(tf, clean(h2.get_text()), size=28, bold=True, color=WHITE, font='Georgia')

    if bd:
        tf = add_txb(slide, M, Inches(1.28), W - 2*M, Inches(5.9))
        first = True
        for item_type, text in extract_items(bd):
            sb = Pt(0) if first else Pt(6)
            if item_type == 'para':
                para(tf, text, size=15, color=GRAY_TXT, space_before=sb, space_after=Pt(4))
            elif item_type == 'bullet':
                para(tf, '•  ' + text, size=14, color=DARK, space_before=Pt(2), space_after=Pt(2))
            elif item_type in ('section-head', 'card-head'):
                para(tf, text, size=14, bold=True, color=NAVY, space_before=sb, space_after=Pt(2))
            elif item_type == 'card-body':
                para(tf, text, size=13, color=GRAY_TXT, space_before=Pt(1), space_after=Pt(5))
            elif item_type == 'feature-label':
                para(tf, text.upper(), size=10, bold=True, color=LT_BLUE, space_before=sb, space_after=Pt(1))
            elif item_type == 'feature-text':
                para(tf, text, size=13, color=DARK, space_before=Pt(0), space_after=Pt(4))
            elif item_type == 'callout-label':
                para(tf, text.upper(), size=10, bold=True, color=MED_BLUE, space_before=sb, space_after=Pt(1))
            elif item_type == 'callout-text':
                para(tf, text, size=13, italic=True, color=DARK, space_before=Pt(0), space_after=Pt(4))
            first = False


# ── Main ──────────────────────────────────────────

def convert(html_path, pptx_path):
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slides = soup.find_all('div', class_='slide')
    print(f"  {len(slides)} slides  ←  {os.path.basename(html_path)}")

    for i, s in enumerate(slides):
        cls = s.get('class', [])
        if 's-title' in cls:
            build_title(prs, s)
        elif 's-divider' in cls:
            build_divider(prs, s)
        elif 's-nutshell' in cls:
            build_nutshell(prs, s)
        else:
            build_std(prs, s)
        print(f"    [{i+1:02d}] {' '.join(cls)}")

    prs.save(pptx_path)
    print(f"  → {os.path.basename(pptx_path)}\n")


if __name__ == '__main__':
    BASE = os.path.dirname(os.path.abspath(__file__))
    modules = [
        'module-1-1-canadian-payments-law',
        'module-2-clearing-settlement-liquidity',
        'module-3-operational-risks-outage-mitigation',
        'module-4-participation-models-credit-unions',
    ]
    for m in modules:
        convert(
            os.path.join(BASE, m + '.html'),
            os.path.join(BASE, m + '.pptx'),
        )
    print("Done — 4 PPTX files ready.")
